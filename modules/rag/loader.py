import logging
from pathlib import Path


class DocumentLoadError(RuntimeError):
    """Raised when a supported document cannot provide usable text for RAG."""


TEXT_SUFFIXES = {".txt", ".md", ".markdown"}


def load_document(path: str | Path) -> str:
    """Load plain text from a supported course document."""
    document_path = Path(path)
    suffix = document_path.suffix.lower()

    if not document_path.exists():
        raise DocumentLoadError(f"资料文件不存在：{document_path}")

    if suffix in TEXT_SUFFIXES:
        try:
            return _require_text(
                document_path.read_text(encoding="utf-8"),
                document_name=document_path.name,
            )
        except UnicodeDecodeError as exc:
            raise DocumentLoadError(
                f"文本文件解析失败：{document_path.name}。请确认文件使用 UTF-8 编码。"
            ) from exc

    if suffix == ".pdf":
        return _load_pdf(document_path)

    if suffix == ".pptx":
        return _load_pptx(document_path)

    raise ValueError(f"Unsupported document type: {suffix}")


def _load_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DocumentLoadError("读取 PDF 需要安装 pypdf，请先安装 requirements.txt。") from exc

    try:
        previous_disable_level = logging.root.manager.disable
        # Broken PDFs can make pypdf emit low-level parser logs before raising;
        # users should see our concise parsing error instead of library noise.
        logging.disable(logging.CRITICAL)
        try:
            reader = PdfReader(str(path))
            page_texts = []
            for page_number, page in enumerate(reader.pages, start=1):
                text = page.extract_text() or ""
                if text.strip():
                    page_texts.append(f"[第 {page_number} 页]\n{text.strip()}")
        finally:
            logging.disable(previous_disable_level)
    except Exception as exc:
        raise DocumentLoadError(
            f"PDF 文件解析失败：{path.name}。请确认文件未损坏；如果是扫描版 PDF，需要先进行 OCR。"
        ) from exc

    return _require_text("\n\n".join(page_texts), document_name=path.name)


def _load_pptx(path: Path) -> str:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocumentLoadError("读取 PPTX 需要安装 python-pptx，请先安装 requirements.txt。") from exc

    try:
        presentation = Presentation(str(path))
        slide_texts = []
        for slide_number, slide in enumerate(presentation.slides, start=1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(f"[第 {slide_number} 页幻灯片]\n" + "\n".join(texts))
    except Exception as exc:
        raise DocumentLoadError(
            f"PPTX 文件解析失败：{path.name}。请确认文件未损坏且格式为 .pptx。"
        ) from exc

    return _require_text("\n\n".join(slide_texts), document_name=path.name)


def _require_text(text: str, document_name: str) -> str:
    # RAG indexing must fail fast on empty extraction; otherwise users see a
    # misleading "indexed" state while every later retrieval returns nothing.
    normalized = text.strip()
    if not normalized:
        raise DocumentLoadError(f"资料解析失败：{document_name} 没有提取到可用文字。")
    return normalized
